#!/usr/bin/env python3
"""Merge new slides into existing presentation"""

from pptx import Presentation
import sys

def merge_presentations(base_pptx, new_slides_pptx, output_pptx):
    """Append slides from new_slides_pptx to base_pptx"""

    # Load presentations
    base_prs = Presentation(base_pptx)
    new_prs = Presentation(new_slides_pptx)

    # Copy slides from new presentation to base
    for slide in new_prs.slides:
        # Get the slide layout and add new slide
        slide_layout = base_prs.slide_layouts[slide.slide_layout.name] if hasattr(slide.slide_layout, 'name') else base_prs.slide_layouts[0]
        new_slide = base_prs.slides.add_slide(slide_layout)

        # Copy all shapes from source slide to new slide
        for shape in slide.shapes:
            el = shape.element
            newel = type(el)(el)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    # Save merged presentation
    base_prs.save(output_pptx)
    print(f"✅ Merged presentation saved: {output_pptx}")
    print(f"   Total slides: {len(base_prs.slides)}")

if __name__ == '__main__':
    base = '/Users/mengblom/Documents/ghx-dex/04-Projects/Exchange_Modernization_Report_May_2026/Presentations/Exchange-Modernization-Update-v4.pptx'
    new_slides = '/Users/mengblom/Documents/ghx-dex/04-Projects/Exchange_Modernization_Report_May_2026/Source/v4/new-slides-11-12.pptx'
    output = '/Users/mengblom/Documents/ghx-dex/04-Projects/Exchange_Modernization_Report_May_2026/Presentations/Exchange-Modernization-Update-v5.pptx'

    merge_presentations(base, new_slides, output)
